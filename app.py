# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
import google.generativeai as genai
from google.api_core import exceptions
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import io
import configparser
import os
import signal
import time
import uuid
from supabase import create_client, Client

# --- 1. APIキー & Supabase設定の取得 ---
def load_api_key():
    config = configparser.ConfigParser()
    file_path = 'APIKEY.ini'
    if os.path.exists(file_path):
        try:
            config.read(file_path, encoding='utf-8-sig')
            return config.get('GEMINI', 'API_KEY')
        except:
            pass
            
    try:
        if "GEMINI" in st.secrets and "API_KEY" in st.secrets["GEMINI"]:
            return st.secrets["GEMINI"]["API_KEY"]
        elif "API_KEY" in st.secrets:
            return st.secrets["API_KEY"]
    except:
        return None
    return None

EMBEDDED_API_KEY = load_api_key()

# Supabase接続情報の取得
def get_supabase_client() -> Client:
    try:
        if "SUPABASE" in st.secrets:
            url = st.secrets["SUPABASE"]["URL"]
            key = st.secrets["SUPABASE"]["KEY"]
            return create_client(url, key)
    except Exception as e:
        st.error(f"Supabase接続エラー: {e}")
    return None

# --- 2. 各ファイル抽出関数 ---
def extract_from_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_from_pdf(file):
    reader = PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_from_pptx(file):
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_from_excel(file):
    all_sheets = pd.read_excel(file, sheet_name=None)
    text_data = []
    for sheet_name, df in all_sheets.items():
        text_data.append(f"--- シート名: {sheet_name} ---\n{df.to_string(index=False)}")
    return "\n".join(text_data)

def create_excel_download(text):
    output = io.BytesIO()
    lines = text.split('\n')
    table_data = []
    for line in lines:
        if '|' in line:
            cells = [c.strip() for c in line.split('|') if c.strip()]
            if cells and not all(c == '-' or c.startswith('---') for c in cells):
                table_data.append(cells)
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        if table_data:
            df = pd.DataFrame(table_data)
            df.to_excel(writer, index=False, header=False, sheet_name='評価結果')
            workbook = writer.book
            worksheet = writer.sheets['評価結果']
            border_fmt = workbook.add_format({'border': 1, 'text_wrap': True, 'valign': 'top'})
            for row_num in range(len(table_data)):
                worksheet.set_row(row_num, None, border_fmt)
            worksheet.set_column(0, 5, 30)
        else:
            df = pd.DataFrame([text.split('\n')])
            df.to_excel(writer, index=False, header=False, sheet_name='評価結果')
    return output.getvalue()

def get_safe_model_name(api_key):
    try:
        genai.configure(api_key=api_key)
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        has_flash = any('gemini-1.5-flash' in m for m in available_models)
        return 'gemini-1.5-flash' if has_flash else available_models[0].replace('models/', '')
    except:
        return 'gemini-1.5-flash'

def generate_initial_greeting(persona, api_key):
    target_model = get_safe_model_name(api_key)
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(target_model)
        prompt = f"あなたは以下の顧客ペルソナとして手続きに来た本人です。最初の第一声を1つだけ作成してください。\n【ペルソナ】\n氏名: {persona['name']}, 年齢: {persona['age']}, 目的: {persona['purpose_situation']}"
        response = model.generate_content(prompt)
        return response.text.strip()
    except:
        return "こんにちは。定期預金の満期の件で来たんですが、今の時代、普通に預けていても全然増えないですよね…"

def get_ai_roleplay_response(messages, persona, product_docs, api_key):
    target_model = get_safe_model_name(api_key)
    recent_messages = [messages[0]] + messages[-5:] if len(messages) > 6 else messages
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(target_model)
        history_text = ""
        for m in recent_messages:
            role_label = "顧客" if m["role"] == "assistant" else "銀行員"
            history_text += f"{role_label}: {m['content']}\n"
        system_prompt = f"あなたは顧客ペルソナになりきり銀行員と対話するAIです。短く返答してください。\n【ペルソナ】\n氏名: {persona['name']}, 設定: {persona['purpose_situation']}\n【履歴】\n{history_text}"
        response = model.generate_content(system_prompt)
        return response.text
    except Exception as e:
        return f"【システムエラー】: {str(e)}"

def generate_evaluation_report(messages, persona, api_key):
    target_model = get_safe_model_name(api_key)
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(target_model)
        history_text = ""
        for m in messages:
            role_label = "顧客" if m["role"] == "assistant" else "銀行員"
            history_text += f"{role_label}: {m['content']}\n"
        prompt = f"以下の会話ログを厳密に評価しレポートをマークダウン形式で作成してください。\n{history_text}"
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"評価レポートの生成に失敗しました: {str(e)}"

# --- 3. Supabaseへのデータ保存関数 ---
def save_chat_to_supabase(session_id, user_input, ai_response):
    supabase = get_supabase_client()
    if supabase:
        try:
            supabase.table("chat_logs").insert({
                "session_id": session_id,
                "user_inputs": user_input,
                "ai_responses": ai_response
            }).execute()
        except Exception as e:
            st.warning(f"ログ保存失敗: {e}")

def save_report_to_supabase(session_id, report_content):
    supabase = get_supabase_client()
    if supabase:
        try:
            supabase.table("evaluation_reports").insert({
                "session_id": session_id,
                "report_content": report_content
            }).execute()
        except Exception as e:
            st.warning(f"レポート保存失敗: {e}")

# --- 4. 画面構築 (Streamlit UI) ---
st.set_page_config(page_title="銀行員向け 金融商品販売AIロールプレイ", layout="wide")
st.title("🏦 金融商品販売 AIロールプレイシステム")

if "session_id" not in st.session_state:
    st.session_state.session_id = str(uuid.uuid4())

st.sidebar.header("👤 顧客ペルソナ設定")
p_name = st.sidebar.text_input("氏名", value="山田 規子")
p_age = st.sidebar.text_input("年齢", value="65歳")
p_job = st.sidebar.text_input("職業", value="専業主婦")
p_family = st.sidebar.text_input("家族構成", value="夫と二人暮らし")
p_purpose_situation = st.sidebar.text_area("来店目的", value="定期預金が満期を迎えたため手続きに来店。")
p_personality = st.sidebar.text_input("性格", value="慎重派で心配性。")
p_experience = st.sidebar.selectbox("投資経験", ["全くない", "少しある", "豊富にある"])

current_persona = {
    "name": p_name, "age": p_age, "job": p_job, "family": p_family,
    "purpose_situation": p_purpose_situation, "personality": p_personality, "experience": p_experience
}

if st.sidebar.button("⚙️ 顧客ペルソナ設定・リセット"):
    st.session_state.persona = current_persona
    st.session_state.session_id = str(uuid.uuid4()) # 新しいセッション
    dynamic_initial_message = generate_initial_greeting(current_persona, EMBEDDED_API_KEY)
    st.session_state.messages = [{"role": "assistant", "content": dynamic_initial_message}]
    st.session_state.report = ""
    st.rerun()

if "persona" not in st.session_state:
    st.session_state.persona = current_persona
if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "assistant", "content": "こんにちは。定期預金の満期の件で来たんですが、今の時代増えないですよね…"}]
if "report" not in st.session_state:
    st.session_state.report = ""

col_chat, col_report = st.columns([1.2, 1.0])

with col_chat:
    st.subheader("💬 ロールプレイ対話画面")
    for m in st.session_state.messages:
        role = "assistant" if m["role"] == "assistant" else "user"
        avatar = "👤" if role == "assistant" else "💼"
        with st.chat_message(role, avatar=avatar):
            st.markdown(m["content"])

    if prompt := st.chat_input("銀行員としての返答・提案を入力してください"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user", avatar="💼"):
            st.markdown(prompt)

        with st.chat_message("assistant", avatar="👤"):
            with st.spinner("考え中..."):
                res = get_ai_roleplay_response(st.session_state.messages, st.session_state.persona, [], EMBEDDED_API_KEY)
            st.markdown(res)
        st.session_state.messages.append({"role": "assistant", "content": res})
        
        # 会話が成立するたびにSupabaseへ自動リアルタイム保存
        save_chat_to_supabase(st.session_state.session_id, prompt, res)
        st.rerun()

with col_report:
    st.subheader("📊 応対評価レポート")
    if st.button("📝 応対評価レポートを生成する", type="primary"):
        if len(st.session_state.messages) <= 1:
            st.warning("会話が開始されていません。")
        else:
            with st.spinner("分析中..."):
                report_res = generate_evaluation_report(st.session_state.messages, st.session_state.persona, EMBEDDED_API_KEY)
                st.session_state.report = report_res
                # レポート生成時にSupabaseへ自動保存
                save_report_to_supabase(st.session_state.session_id, report_res)
                st.rerun()
                
    if st.session_state.report:
        st.markdown(st.session_state.report)

